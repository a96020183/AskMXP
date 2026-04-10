"""Build the 3-page AskMXP 2.0 pitch deck as a standalone .pptx."""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


NAVY = RGBColor(0x0F, 0x1E, 0x3D)
DEEP_BLUE = RGBColor(0x1B, 0x2A, 0x4E)
ACCENT_RED = RGBColor(0xE6, 0x21, 0x17)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)
ACCENT_GOLD = RGBColor(0xFF, 0xC3, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY = RGBColor(0x22, 0x22, 0x22)
SOFT_BLUE = RGBColor(0x2B, 0x48, 0x7D)
CARD_DARK = RGBColor(0x16, 0x27, 0x4A)

FONT_ZH = "Microsoft JhengHei"


def add_bg(slide, prs, color=NAVY):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name=FONT_ZH):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def add_rect(slide, left, top, width, height, fill,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08):
    r = slide.shapes.add_shape(shape, left, top, width, height)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    r.line.fill.background()
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            r.adjustments[0] = radius
        except Exception:
            pass
    return r


def add_line(slide, left, top, width, height, color):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_bottom_banner(slide, prs, text, color=DEEP_BLUE, text_color=ACCENT_GOLD):
    banner = add_rect(slide,
                      0, prs.slide_height - Inches(0.7),
                      prs.slide_width, Inches(0.7),
                      color, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, 0, prs.slide_height - Inches(0.62),
             prs.slide_width, Inches(0.55),
             text, size=18, bold=True, color=text_color,
             align=PP_ALIGN.CENTER)
    return banner


def add_page_tag(slide, prs, tag_text):
    add_text(slide,
             prs.slide_width - Inches(1.4), Inches(0.35),
             Inches(1.2), Inches(0.4),
             tag_text, size=12, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ─────────────── SLIDE 1 ───────────────
def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, NAVY)

    # left accent bar
    add_line(slide, 0, 0, Inches(0.25), prs.slide_height, ACCENT_RED)

    add_page_tag(slide, prs, "01 / 03")

    add_text(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.5),
             "PRODUCT VISION", size=12, bold=True, color=ACCENT_ORANGE)
    add_text(slide, Inches(0.7), Inches(0.95), Inches(12.3), Inches(1.1),
             "從被動看數據，到主動抓洞察", size=40, bold=True, color=WHITE)
    add_text(slide, Inches(0.7), Inches(1.95), Inches(12), Inches(0.5),
             "AskMXP 2.0 · 智慧數據監控系統", size=18, color=LIGHT_GRAY)

    # Pain card 1
    add_rect(slide, Inches(0.7), Inches(2.75), Inches(5.8), Inches(1.9), ACCENT_RED)
    add_text(slide, Inches(0.9), Inches(2.85), Inches(5.5), Inches(1.1),
             "2 小時 / 天", size=48, bold=True, color=WHITE)
    add_text(slide, Inches(0.9), Inches(4.0), Inches(5.5), Inches(0.55),
             "PM 每日人工巡檢成本", size=15, color=WHITE)

    # Pain card 2
    add_rect(slide, Inches(0.7), Inches(4.85), Inches(5.8), Inches(1.9), ACCENT_ORANGE)
    add_text(slide, Inches(0.9), Inches(4.95), Inches(5.5), Inches(1.1),
             "20% 跌幅漏接", size=48, bold=True, color=WHITE)
    add_text(slide, Inches(0.9), Inches(6.1), Inches(5.5), Inches(0.55),
             "傳統人工監控失效率", size=15, color=WHITE)

    # Before / After panel (right)
    add_rect(slide, Inches(6.9), Inches(2.75), Inches(6.1), Inches(4.0), CARD_DARK)

    add_text(slide, Inches(7.2), Inches(2.9), Inches(5.5), Inches(0.4),
             "▍ BEFORE", size=13, bold=True, color=MID_GRAY)
    add_text(slide, Inches(7.2), Inches(3.3), Inches(5.5), Inches(1.3),
             "PM 疲於切換多個 Dashboard\n手動複製數字到 Excel\n異常跌幅常在切換之間漏看",
             size=15, color=WHITE)

    add_line(slide, Inches(7.2), Inches(4.75), Inches(5.5), Inches(0.03), ACCENT_RED)

    add_text(slide, Inches(7.2), Inches(4.9), Inches(5.5), Inches(0.4),
             "▍ AFTER", size=13, bold=True, color=ACCENT_GOLD)
    add_text(slide, Inches(7.2), Inches(5.3), Inches(5.5), Inches(1.3),
             "Google Chat 自動跳出結構化警示卡\n24/7 主動監控，跌幅 > 20% 秒級通知\n異常會自己找上你，而不是你去找它",
             size=15, color=WHITE)

    add_bottom_banner(slide, prs, "把數據監控，從 PM 的義務變成系統的防禦。")

    set_notes(slide, (
        "大家好，我是 APM 實習生。先丟一個問題：你知道我們 PM 每天花多少時間在"
        "『看數據』嗎？答案是平均 2 小時。而且更殘酷的是，根據我過去一個月的"
        "實地觀察，大約有 20% 的異常跌幅會在人工巡檢時被漏掉。\n\n"
        "為什麼？因為數據分散在 Mixpanel、GA、後台各處，而 PM 的注意力是有限的。"
        "當我們把時間花在切換分頁、複製數字、貼到 Excel，真正該做產品決策的時間"
        "就被壓縮了。\n\n"
        "所以我想解決的不是『做一個更漂亮的 Dashboard』，而是要讓數據監控這件事，"
        "從『每天的義務』變成『24 小時的自動防禦』。這就是 AskMXP 2.0 的起點。"
    ))


# ─────────────── SLIDE 2 ───────────────
def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, NAVY)
    add_line(slide, 0, 0, Inches(0.25), prs.slide_height, ACCENT_ORANGE)

    add_page_tag(slide, prs, "02 / 03")

    add_text(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.5),
             "CORE FEATURES", size=12, bold=True, color=ACCENT_ORANGE)
    add_text(slide, Inches(0.7), Inches(0.95), Inches(12.3), Inches(1.1),
             "雙引擎核心：機器巡檢 × AI 記憶", size=36, bold=True, color=WHITE)
    add_text(slide, Inches(0.7), Inches(1.95), Inches(12), Inches(0.5),
             "亮點 A 自動化警示    ×    亮點 B AI 因果分析",
             size=16, color=LIGHT_GRAY)

    # Left column — Highlight A
    add_rect(slide, Inches(0.7), Inches(2.75), Inches(6.0), Inches(4.0), CARD_DARK)
    add_text(slide, Inches(0.9), Inches(2.85), Inches(5.6), Inches(0.45),
             "亮點 A · 主動防禦", size=13, bold=True, color=ACCENT_ORANGE)
    add_text(slide, Inches(0.9), Inches(3.25), Inches(5.6), Inches(0.6),
             "GitHub Actions 自動化監控", size=22, bold=True, color=WHITE)

    add_text(slide, Inches(0.9), Inches(3.95), Inches(5.6), Inches(0.5),
             "⏰  每日 09:00 自動執行", size=14, color=LIGHT_GRAY)

    add_rect(slide, Inches(0.9), Inches(4.55), Inches(5.6), Inches(0.75), SOFT_BLUE)
    add_text(slide, Inches(1.0), Inches(4.65), Inches(5.4), Inches(0.55),
             "Mixpanel → 7 日均值比對 → Chat 警示",
             size=13, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0.9), Inches(5.5), Inches(5.6), Inches(1.2),
             "■ 跌幅超過 20% 自動推卡片\n"
             "■ 24/7 監控，無需人力值守\n"
             "■ 從「被動巡檢」→「主動防禦」",
             size=13, color=WHITE)

    # Right column — Highlight B
    add_rect(slide, Inches(7.0), Inches(2.75), Inches(6.0), Inches(4.0), CARD_DARK)
    add_text(slide, Inches(7.2), Inches(2.85), Inches(5.6), Inches(0.45),
             "亮點 B · 因果分析", size=13, bold=True, color=ACCENT_GOLD)
    add_text(slide, Inches(7.2), Inches(3.25), Inches(5.6), Inches(0.6),
             "AI 營銷備忘錄", size=22, bold=True, color=WHITE)

    add_text(slide, Inches(7.2), Inches(3.95), Inches(5.6), Inches(0.5),
             "🧠  Claude Sonnet 4.6 深度分析", size=14, color=LIGHT_GRAY)

    add_rect(slide, Inches(7.2), Inches(4.55), Inches(5.6), Inches(0.75), SOFT_BLUE)
    add_text(slide, Inches(7.3), Inches(4.65), Inches(5.4), Inches(0.55),
             "KOL 活動 + 改版紀錄 + Mixpanel → AI",
             size=13, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(7.2), Inches(5.5), Inches(5.6), Inches(1.2),
             "■ 支援歷史活動標註\n"
             "■ 交叉比對數據波動 × 營銷事件\n"
             "■ 從「發生了什麼」→「為什麼發生」",
             size=13, color=WHITE)

    add_bottom_banner(slide, prs, "不告訴你『發生了什麼』,而是告訴你『為什麼發生』。")

    set_notes(slide, (
        "AskMXP 2.0 的核心有兩顆引擎。\n\n"
        "第一顆是『自動化警示』。我把監控邏輯部署在 GitHub Actions 上，每天早上"
        "九點自動抓 Mixpanel 資料，把今天的數值和過去七天的均值做比對，只要跌幅"
        "超過 20%，就會直接推一張結構化的警示卡到 Google Chat。這代表 PM 不用"
        "再主動去看，而是異常會自己找上你——這是從『被動巡檢』變成『主動防禦』。\n\n"
        "第二顆是『AI 營銷記憶』。單純看數字永遠只能說『漲了、跌了』，但 PM 真正"
        "想知道的是『為什麼』。所以我讓使用者可以把 KOL 活動、改版紀錄、推播時間"
        "標註進系統，當 AI 產出備忘錄時，它會把這些歷史活動和數據波動做交叉比對，"
        "直接告訴你『這個尖峰很可能來自 6/7 的 KOL 合作』。這就是把 Mixpanel 從"
        "一個『數據倉庫』升級成『一個會思考的分析師』。"
    ))


# ─────────────── SLIDE 3 ───────────────
def build_slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, NAVY)
    add_line(slide, 0, 0, Inches(0.25), prs.slide_height, ACCENT_GOLD)

    add_page_tag(slide, prs, "03 / 03")

    add_text(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.5),
             "EFFICIENCY REVOLUTION", size=12, bold=True, color=ACCENT_ORANGE)
    add_text(slide, Inches(0.7), Inches(0.95), Inches(12.3), Inches(1.1),
             "從 1 小時剪貼，到 10 秒一鍵出報", size=36, bold=True, color=WHITE)
    add_text(slide, Inches(0.7), Inches(1.95), Inches(12), Inches(0.5),
             "亮點 C · 一鍵自動化週報產出",
             size=16, color=LIGHT_GRAY)

    # Before bar
    add_text(slide, Inches(0.7), Inches(2.8), Inches(1.5), Inches(0.4),
             "原本", size=14, bold=True, color=LIGHT_GRAY)
    add_rect(slide, Inches(2.3), Inches(2.75), Inches(7.5), Inches(0.6),
             MID_GRAY, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, Inches(2.5), Inches(2.82), Inches(7.2), Inches(0.45),
             "1 小時｜手動截圖 + 貼 PPT + 寫摘要",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(10.0), Inches(2.8), Inches(3.0), Inches(0.4),
             "60 min", size=16, bold=True, color=LIGHT_GRAY)

    # After bar
    add_text(slide, Inches(0.7), Inches(3.75), Inches(1.5), Inches(0.4),
             "現在", size=14, bold=True, color=ACCENT_GOLD)
    add_rect(slide, Inches(2.3), Inches(3.7), Inches(0.7), Inches(0.6),
             ACCENT_RED, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, Inches(2.35), Inches(3.77), Inches(0.6), Inches(0.45),
             "10 sec", size=10, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.1), Inches(3.78), Inches(6.5), Inches(0.4),
             "一鍵點擊 · AskMXP 自動生成 .pptx",
             size=14, bold=True, color=WHITE)
    add_text(slide, Inches(10.0), Inches(3.75), Inches(3.0), Inches(0.4),
             "× 360 倍效率", size=16, bold=True, color=ACCENT_GOLD)

    # PPT output spec
    add_rect(slide, Inches(0.7), Inches(4.6), Inches(7.8), Inches(2.15), CARD_DARK)
    add_text(slide, Inches(0.9), Inches(4.72), Inches(7.4), Inches(0.45),
             "AskMXP 週報內容自動產出", size=13, bold=True, color=ACCENT_ORANGE)
    add_text(slide, Inches(0.9), Inches(5.1), Inches(7.4), Inches(1.6),
             "📊  每事件一張投影片，30 天趨勢折線圖\n"
             "🤖  Claude 撰寫的專業營銷洞察文字\n"
             "🏷️  自動比對已標註的歷史活動\n"
             "⬇️  直接下載 .pptx，不用再動剪刀手",
             size=14, color=WHITE)

    # Conclusion block
    add_rect(slide, Inches(8.8), Inches(4.6), Inches(4.2), Inches(2.15), ACCENT_RED)
    add_text(slide, Inches(9.0), Inches(4.85), Inches(3.9), Inches(0.9),
             "把時間\n還給 PM", size=32, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_line(slide, Inches(9.8), Inches(6.05), Inches(2.2), Inches(0.03), WHITE)
    add_text(slide, Inches(9.0), Inches(6.1), Inches(3.9), Inches(0.8),
             "讓人思考策略\n而不是複製貼上數據",
             size=13, color=WHITE, align=PP_ALIGN.CENTER)

    add_bottom_banner(slide, prs, "AskMXP 2.0 · Thank You",
                      color=DEEP_BLUE, text_color=WHITE)

    set_notes(slide, (
        "第三個亮點，也是我自己最有感的一個——週報自動化。過去產出一份有圖表的"
        "週報，PM 至少要花 1 個小時在 Mixpanel 截圖、貼進 PPT、手動寫摘要。"
        "現在，AskMXP 只要按一個按鈕，10 秒內就會產出一份完整的 pptx 檔案，"
        "裡面每一頁都是一個事件，包含 30 天趨勢圖、AI 撰寫的營銷洞察，以及對應的"
        "歷史活動標註。效率提升了 360 倍。\n\n"
        "最後我想說，這個工具真正想解決的不是『出報告』這件事，而是"
        "『PM 的時間被什麼佔滿』這件事。當監控自動化、分析 AI 化、報告一鍵化，"
        "PM 省下來的不是 3 小時，而是可以拿來思考策略的 3 小時。\n\n"
        "AskMXP 2.0 的初衷就是：把時間還給 PM，讓人去做人該做的事——"
        "思考策略、洞察用戶、定義產品。謝謝大家。"
    ))


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)

    out = "AskMXP_Pitch_3pages.pptx"
    prs.save(out)
    print(f"✅ 已產出：{out}")


if __name__ == "__main__":
    main()
