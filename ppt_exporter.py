"""Weekly report PPT exporter using python-pptx + kaleido."""
from __future__ import annotations

import io
from datetime import datetime
from typing import Iterable

import pandas as pd
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _make_chart_png(df: pd.DataFrame, title: str) -> bytes:
    fig = go.Figure()
    fig.add_trace(go.Scatter(
        x=df["日期"], y=df["事件次數"],
        mode="lines+markers",
        line=dict(color="#E62117", width=2),
        marker=dict(size=5),
    ))
    fig.update_layout(
        title=title,
        xaxis_title="日期",
        yaxis_title="事件次數",
        template="plotly_white",
        width=900,
        height=450,
        margin=dict(l=50, r=30, t=60, b=50),
    )
    return fig.to_image(format="png", engine="kaleido")


def _add_cover(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _add_event_slide(prs: Presentation, chinese_name: str, event_id: str,
                     df: pd.DataFrame, memo: str) -> None:
    layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = f"{chinese_name}"

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.4))
    tf = subtitle_box.text_frame
    tf.text = f"事件 ID：{event_id}"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x77, 0x77, 0x77)

    if not df.empty:
        png = _make_chart_png(df, f"{chinese_name} 過去 30 天每日觸發次數")
        image_stream = io.BytesIO(png)
        slide.shapes.add_picture(
            image_stream,
            Inches(0.5), Inches(1.7),
            width=Inches(7.5), height=Inches(3.75),
        )

    memo_box = slide.shapes.add_textbox(Inches(8.2), Inches(1.7), Inches(4.8), Inches(5.4))
    tf = memo_box.text_frame
    tf.word_wrap = True
    lines = memo.strip().split("\n") or ["(無摘要)"]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)


def build_weekly_report(events_results: Iterable[dict], output_path: str) -> str:
    """events_results: iterable of dicts with keys
       {chinese_name, event_id, df, memo}."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    today = datetime.now().strftime("%Y-%m-%d")
    _add_cover(prs, "AskMXP 週報", f"Mixpanel 事件監控 · {today}")

    count = 0
    for ev in events_results:
        _add_event_slide(
            prs,
            chinese_name=ev["chinese_name"],
            event_id=ev["event_id"],
            df=ev["df"],
            memo=ev.get("memo", ""),
        )
        count += 1

    summary_slide = prs.slides.add_slide(prs.slide_layouts[5])
    summary_slide.shapes.title.text = "報告結尾"
    end_box = summary_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(4))
    end_box.text_frame.text = f"本週共分析 {count} 個事件。\n產生時間：{today}"

    prs.save(output_path)
    return output_path
